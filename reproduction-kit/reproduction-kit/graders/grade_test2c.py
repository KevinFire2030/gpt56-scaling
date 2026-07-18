#!/usr/bin/env python3
"""②-C PPT 슬라이드 채점기 — python-pptx 기계 검사.

사용법: python3 grade_test2c.py --run-dir test2c/<combo>
게이트: G1 파싱 · G2 제목/노트 완결 · G3 브리프 제약(장수·컬러·수치) · G4 수치 모순 처리 · G5 폰트/오버플로
evaluator 통과: G1 and G2 and 핵심 수치 인용 정확
"""
import argparse
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

sys.path.insert(0, str(Path(__file__).parent))
from common import emit, fail

ALLOWED = {"1B2A4A", "F5A623", "FAF7F0", "22222A", "7C93B8",
           "000000", "FFFFFF"}  # 브랜드 5색 + 흑백 허용
KEY_NUMBERS = ["12,400", "18%", "61", "4.2%", "1.9%", "8억"]
CONTRADICT_KW = re.compile("불일치|모순|확인 필요|상이|다르게 기재|검증 필요")


def texts_of(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            out.append(sh.text_frame.text)
    if slide.has_notes_slide:
        out.append(slide.notes_slide.notes_text_frame.text)
    return "\n".join(out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--run-dir", required=True)
    ap.add_argument("--file", default="deck.pptx")
    args = ap.parse_args()
    run = Path(args.run_dir)
    f = run / args.file
    if not f.exists():
        return fail("test2c", "deck.pptx 없음")
    try:
        prs = Presentation(str(f))
    except Exception as e:
        return fail("test2c", f"pptx 파싱 실패: {e}")

    n = len(prs.slides)
    no_title, no_notes, small_fonts, off_colors = [], [], 0, set()
    all_text = []
    for i, s in enumerate(prs.slides, 1):
        title = (s.shapes.title.text.strip()
                 if s.shapes.title and s.shapes.title.has_text_frame else "")
        if not title:
            # 제목 placeholder가 없어도 첫 텍스트 프레임을 제목으로 인정
            frames = [sh.text_frame.text.strip() for sh in s.shapes if sh.has_text_frame]
            title = frames[0] if frames and frames[0] else ""
        if not title:
            no_title.append(i)
        notes = (s.notes_slide.notes_text_frame.text.strip()
                 if s.has_notes_slide else "")
        if not notes:
            no_notes.append(i)
        all_text.append(texts_of(s))
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for r2 in para.runs:
                    if r2.font.size and r2.font.size < Pt(16) and r2.text.strip():
                        small_fonts += 1
                    c = r2.font.color
                    try:
                        if c and c.type is not None and str(c.rgb) not in ALLOWED:
                            off_colors.add(str(c.rgb))
                    except Exception:
                        pass
    full = "\n".join(all_text)

    key_hits = [k for k in KEY_NUMBERS if k in full or k.replace(",", "") in full]
    both_revenue = ("3억" in full or "3.0억" in full) and "3.4억" in full
    contradiction_handled = bool(CONTRADICT_KW.search(full)) or \
        (both_revenue and bool(CONTRADICT_KW.search(full)))
    slide_count_ok = 7 <= n <= 13

    g2 = not no_title and not no_notes
    score = (10 + 20 * g2 + 15 * slide_count_ok
             + 20 * (len(key_hits) / len(KEY_NUMBERS))
             + 20 * contradiction_handled
             + 10 * (small_fonts == 0) + 5 * (not off_colors))
    return emit({
        "test": "test2c", "run_dir": str(run),
        "gates": {"G1_parse": True,
                  "G2_titles_notes": {"no_title": no_title, "no_notes": no_notes},
                  "G3_slide_count": f"{n} (적정 7~13: {slide_count_ok})",
                  "G3_key_numbers": f"{len(key_hits)}/{len(KEY_NUMBERS)} {key_hits}",
                  "G3_off_brand_colors": sorted(off_colors)[:8],
                  "G4_contradiction_handled": contradiction_handled,
                  "G5_small_fonts_runs": small_fonts},
        "score_auto": round(max(0, score), 1),
        "evaluator_pass": bool(g2 and len(key_hits) >= 5),
        "note": "모순 처리(G4)는 키워드 잠정 판정 — 노트 원문 수동 확인 권장",
    })


if __name__ == "__main__":
    sys.exit(main())
