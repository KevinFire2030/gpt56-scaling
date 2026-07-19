from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pathlib import Path

OUT = Path('deck.pptx')
W, H = 13.333, 7.5
C = {
    'navy': '1B2A4A', 'amber': 'F5A623', 'paper': 'FAF7F0',
    'ink': '22222A', 'steel': '7C93B8', 'white': 'FFFFFF',
    'mist': 'E8EDF3', 'pale_amber': 'FCE7BF', 'green': '3C7A67',
    'red': 'B94A48'
}

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
prs.core_properties.title = 'Aurora Kit 출시 6주 성과 보고'
prs.core_properties.subject = '경영진 의사결정용 신제품 출시 성과 보고'
prs.core_properties.author = '경영진 보고용 자동 생성 자료'
blank = prs.slide_layouts[6]

def rgb(hexstr): return RGBColor.from_string(hexstr)

def rect(slide, x,y,w,h, fill, line=None, radius=False, transparency=0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(C[fill] if fill in C else fill)
    if transparency: shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(C[line] if line and line in C else (line or (C[fill] if fill in C else fill)))
    if not line: shape.line.fill.background()
    return shape

def txt(slide, text, x,y,w,h, size=16, color='ink', bold=False, align=PP_ALIGN.LEFT, font='Aptos', valign=MSO_ANCHOR.TOP, margin=0.03):
    box = slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Inches(margin); tf.margin_top=tf.margin_bottom=Inches(margin)
    tf.vertical_anchor=valign
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(C[color] if color in C else color)
    return box

def line(slide,x1,y1,x2,y2,color='steel',width=1.5,dash=None):
    ln=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    ln.line.color.rgb=rgb(C[color] if color in C else color); ln.line.width=Pt(width)
    if dash: ln.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return ln

def base(slide, dark=False, section='AURORA KIT / 6주 성과 보고'):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb=rgb(C['navy'] if dark else C['paper'])
    txt(slide, section, 0.6, 0.28, 4.8, 0.22, 9, 'white' if dark else 'steel', bold=True, font='Aptos', margin=0)
    txt(slide, f'{len(prs.slides):02d}', 12.15, 7.08, 0.55, 0.2, 9, 'white' if dark else 'steel', align=PP_ALIGN.RIGHT, margin=0)

def title(slide, t, sub=None, dark=False):
    txt(slide,t,0.6,0.72,12.0,0.62,30,'white' if dark else 'navy',bold=True,font='Aptos Display',margin=0)
    if sub: txt(slide,sub,0.62,1.38,11.8,0.35,13,'mist' if dark else 'steel',margin=0)

def notes(slide, content):
    for shp in slide.notes_slide.placeholders:
        if shp.placeholder_format.type == 2:
            shp.text = content
            return

def stat_card(slide,x,y,w,h,value,label,accent='amber',dark=False,small=''):
    rect(slide,x,y,w,h,'white' if not dark else 'mist',line='mist' if not dark else None,radius=True)
    rect(slide,x,y,0.08,h,accent)
    txt(slide,value,x+0.25,y+0.25,w-0.4,0.5,28,'navy',True,margin=0)
    txt(slide,label,x+0.25,y+0.86,w-0.4,0.45,13,'ink' if not dark else 'navy',bold=True,margin=0)
    if small: txt(slide,small,x+0.25,y+1.24,w-0.4,0.25,10,'steel',margin=0)

def add_bullet(slide, text, x,y,w,h, size=17, color='ink', accent='amber'):
    rect(slide,x,y+0.08,0.12,0.12,accent,radius=True)
    txt(slide,text,x+0.28,y,w-0.28,h,size,color,margin=0)

# 1 Cover
s=prs.slides.add_slide(blank); base(s, dark=True)
rect(s,8.65,0,4.68,7.5,'amber')
# abstract Aurora motif
for i,(x,y,w,h,fill) in enumerate([(9.2,1.2,2.8,2.8,'navy'),(10.25,2.3,2.0,2.0,'paper'),(9.65,4.35,2.9,1.15,'navy'),(11.25,5.0,1.2,1.2,'paper')]):
    sh=rect(s,x,y,w,h,fill,radius=True); sh.rotation = [12,-18,8,-22][i]
txt(s,'신제품 출시 성과 보고',0.72,1.25,7.2,0.55,30,'white',True,font='Aptos Display',margin=0)
txt(s,'Aurora Kit · 출시 후 6주',0.75,2.0,6.5,0.55,23,'amber',True,margin=0)
txt(s,'초기 수요는 강하다.\n이제 수익성과 회전을 기준으로 확장한다.',0.76,3.1,6.5,1.25,22,'white',margin=0)
rect(s,0.76,5.55,2.2,0.08,'amber')
txt(s,'경영진 의사결정 요청 포함 · 10분',0.76,5.82,5.0,0.3,12,'mist',margin=0)
notes(s,'오늘은 Aurora Kit 출시 후 6주 성과를 숫자 중심으로 보고드리고, 다음 분기 2차 캠페인의 집행 방향과 추가 마케팅 예산 8억 원 요청을 의사결정 안건으로 올립니다. 핵심은 수요 자체는 강하지만 채널별 수익성과 재고 회전의 편차를 관리하면서 확장해야 한다는 점입니다.')

# 2 agenda
s=prs.slides.add_slide(blank); base(s); title(s,'오늘의 흐름','성과를 확인하고, 실행 기준과 의사결정 안건으로 연결합니다.')
items=[('01','성과 요약','6주 성적표와 목표 대비 결과'),('02','핵심 지표','판매·전환·반품·고객 반응'),('03','실행 시사점','채널 효율과 운영 리스크'),('04','결정 요청','다음 분기 재배분 및 예산')]
for i,(n,h,d) in enumerate(items):
    x=0.8+i*3.05
    rect(s,x,2.35,2.55,2.55,'white',line='mist',radius=True)
    txt(s,n,x+0.22,2.62,0.6,0.35,14,'amber',True,margin=0)
    txt(s,h,x+0.22,3.18,2.1,0.4,20,'navy',True,margin=0)
    txt(s,d,x+0.22,3.82,2.0,0.72,13,'steel',margin=0)
    if i<3: line(s,x+2.62,3.63,x+3.0,3.63,'steel',2)
notes(s,'발표는 네 단계로 진행합니다. 먼저 6주 성과를 한 장으로 요약하고, 핵심 지표를 통해 무엇이 잘됐는지 확인합니다. 이어서 채널 수익성과 재고 회전 관점의 실행 시사점, 데이터 확인 과제를 짚고, 마지막으로 다음 분기 재배분과 예산 요청에 대한 결정을 부탁드립니다.')

# 3 exec summary
s=prs.slides.add_slide(blank); base(s); title(s,'핵심 결론','초기 수요는 예상보다 강합니다. 확장은 선택과 집중이 전제입니다.')
stat_card(s,0.75,2.1,2.75,1.85,'12,400대','출시 후 6주 누적 판매량','amber',small='목표 대비 +18%')
stat_card(s,3.75,2.1,2.75,1.85,'4.2%','웹사이트 방문 대비 구매 전환율','navy',small='초기 유입의 구매 신호')
stat_card(s,6.75,2.1,2.75,1.85,'1.9%','반품률','green',small='품질·기대 정합성 양호')
stat_card(s,9.75,2.1,2.75,1.85,'61개','재구매 의향 조사 응답 기업 고객','steel',small='반응 데이터 축적 중')
rect(s,0.78,4.65,11.6,1.12,'navy',radius=True)
txt(s,'경영진에게 필요한 한 가지 판단',1.05,4.9,3.25,0.26,12,'amber',True,margin=0)
txt(s,'2차 캠페인을 고효율 채널 중심으로 재배분하고, 수익성·재고 회전 기준을 함께 운영한다.',1.05,5.27,10.9,0.32,19,'white',True,margin=0)
notes(s,'성과의 방향은 명확합니다. 6주 누적 판매량은 12,400대이고 목표 대비 판매 초과율은 18%입니다. 구매 전환율 4.2%, 반품률 1.9%로 초기 퍼널과 제품 경험도 긍정적입니다. 다만 2차 캠페인은 모든 채널을 동일하게 키우기보다 고효율 채널 중심으로 재배분하고, 수익성과 재고 회전을 동시에 관리하는 방식이 적절합니다.')

# 4 KPI scorecard
s=prs.slides.add_slide(blank); base(s); title(s,'핵심 지표 성적표','강한 수요와 양호한 제품 신호가 함께 나타났습니다.')
# left large gauge
rect(s,0.75,2.05,4.15,3.95,'white',line='mist',radius=True)
txt(s,'판매 성과',1.08,2.35,2.5,0.3,14,'steel',True,margin=0)
txt(s,'12,400대',1.05,2.9,3.35,0.72,38,'navy',True,margin=0)
txt(s,'출시 후 6주 누적 판매량',1.08,3.72,3.2,0.3,14,'ink',margin=0)
# progress bar
rect(s,1.08,4.45,3.45,0.24,'mist',radius=True); rect(s,1.08,4.45,2.98,0.24,'amber',radius=True)
txt(s,'목표 대비 판매 초과율',1.08,4.92,2.5,0.25,12,'steel',margin=0)
txt(s,'18%',3.9,4.82,0.6,0.35,20,'amber',True,align=PP_ALIGN.RIGHT,margin=0)
# right rows
rows=[('전환율','4.2%','웹사이트 방문 대비 구매'),('반품률','1.9%','초기 품질·기대 정합성'),('재구매 의향','61개 기업 고객','조사 응답 기반')]
for i,(a,b,c) in enumerate(rows):
    y=2.12+i*1.22
    rect(s,5.35,y,6.95,0.92,'white',line='mist',radius=True)
    txt(s,a,5.68,y+0.2,2.0,0.25,13,'steel',True,margin=0)
    txt(s,b,7.65,y+0.13,2.15,0.34,22,'navy',True,margin=0)
    txt(s,c,9.85,y+0.22,2.1,0.28,12,'ink',margin=0)
notes(s,'지표를 세 묶음으로 보겠습니다. 첫째 판매량과 목표 대비 초과율은 수요의 크기를 보여줍니다. 둘째 4.2% 전환율은 유입이 실제 구매로 이어지는 신호입니다. 셋째 반품률 1.9%와 재구매 의향 조사 응답 61개 기업 고객은 제품 경험이 크게 훼손되지 않았음을 시사합니다. 따라서 다음 단계의 과제는 수요 창출보다 효율적인 확장입니다.')

# 5 channel implication
s=prs.slides.add_slide(blank); base(s); title(s,'채널별 시사점','채널을 늘리는 것보다, 효율 기준으로 재배분하는 것이 우선입니다.')
# funnel
labels=[('유입','웹사이트 방문'),('전환','4.2% 구매 전환'),('수익성','채널별 편차'),('회전','재고 회전 편차')]
for i,(h,d) in enumerate(labels):
    x=0.95+i*2.95; w=2.25-0.18*i
    rect(s,x,2.35,w,1.05,'navy' if i<2 else 'amber',radius=True)
    txt(s,h,x+0.15,2.58,w-0.3,0.25,16,'white' if i<2 else 'navy',True,align=PP_ALIGN.CENTER,margin=0)
    txt(s,d,x+0.15,2.95,w-0.3,0.22,11,'mist' if i<2 else 'navy',align=PP_ALIGN.CENTER,margin=0)
    if i<3: line(s,x+w+0.12,2.87,x+w+0.63,2.87,'steel',2)
rect(s,0.95,4.15,11.1,1.05,'white',line='mist',radius=True)
txt(s,'실행 원칙',1.25,4.42,1.4,0.25,13,'amber',True,margin=0)
txt(s,'고효율 채널 중심 재배분  →  채널별 수익성 모니터링  →  재고 회전 기준으로 물량 조정',2.75,4.34,8.9,0.38,17,'navy',True,margin=0)
rect(s,0.95,5.55,11.1,0.45,'pale_amber',radius=True)
txt(s,'주의: 브리프에 채널별 정량 데이터는 없어, 2차 캠페인의 세부 믹스는 성과·수익성 데이터 확인 후 확정',1.15,5.68,10.7,0.18,11,'ink',margin=0)
notes(s,'브리프에 채널별 세부 수치는 없으므로 특정 채널의 성과를 새로 가정하지 않았습니다. 대신 의사결정 원칙을 명확히 제시합니다. 2차 캠페인은 고효율 채널 중심으로 재배분하고, 채널별 수익성을 모니터링하며, 재고 회전 속도에 따라 물량을 조정합니다. 세부 믹스는 채널별 성과와 수익성 데이터를 확인한 뒤 확정하는 것이 안전합니다.')

# 6 customer response
s=prs.slides.add_slide(blank); base(s); title(s,'고객 반응','초기 고객 신호는 긍정적입니다. 다음은 반복 구매로 연결할 단계입니다.')
# big quote style
rect(s,0.8,2.05,5.0,3.75,'navy',radius=True)
txt(s,'“',1.2,2.35,0.6,0.65,54,'amber',True,margin=0)
txt(s,'반품률 1.9%와\n재구매 의향 조사 응답\n61개 기업 고객',1.3,3.05,3.95,1.45,24,'white',True,margin=0)
txt(s,'제품 경험의 초기 안정성',1.3,5.15,3.5,0.25,12,'mist',margin=0)
# right action ladder
for i,(h,d) in enumerate([('1. 신호 확인','61개 기업 고객 응답을 세그먼트별로 읽기'),('2. 반복 구매 설계','재구매 의향이 높은 고객군에 후속 제안'),('3. 학습 확장','반품 사유와 구매 전환 맥락을 캠페인에 반영')]):
    y=2.28+i*1.12
    rect(s,6.35,y,5.45,0.78,'white',line='mist',radius=True)
    rect(s,6.55,y+0.22,0.36,0.36,'amber',radius=True)
    txt(s,str(i+1),6.55,y+0.265,0.36,0.18,12,'navy',True,align=PP_ALIGN.CENTER,margin=0)
    txt(s,h,7.15,y+0.12,2.0,0.22,14,'navy',True,margin=0)
    txt(s,d,7.15,y+0.43,4.2,0.2,11,'ink',margin=0)
notes(s,'고객 반응은 확장에 필요한 기반이 있습니다. 반품률은 1.9%이고, 재구매 의향 조사에는 61개 기업 고객이 응답했습니다. 다음 단계는 이 응답을 세그먼트별로 읽고, 재구매 의향이 높은 고객군에 후속 제안을 설계하는 것입니다. 반품 사유와 구매 전환 맥락도 2차 캠페인의 메시지와 운영에 반영하겠습니다.')

# 7 risks
s=prs.slides.add_slide(blank); base(s); title(s,'리스크 및 확인 과제','현재 가장 명확한 리스크는 매출 기준의 불일치입니다.')
# revenue mismatch
rect(s,0.8,2.05,6.2,3.85,'white',line='mist',radius=True)
txt(s,'누적 매출 기준',1.15,2.35,2.2,0.28,14,'steel',True,margin=0)
rect(s,1.15,3.0,2.25,1.08,'pale_amber',radius=True)
txt(s,'3억 원',1.35,3.25,1.85,0.35,25,'navy',True,align=PP_ALIGN.CENTER,margin=0)
txt(s,'영업팀 주간 보고',1.35,3.67,1.85,0.2,11,'ink',align=PP_ALIGN.CENTER,margin=0)
rect(s,4.2,3.0,2.25,1.08,'mist',radius=True)
txt(s,'3.4억 원',4.4,3.25,1.85,0.35,25,'navy',True,align=PP_ALIGN.CENTER,margin=0)
txt(s,'재무팀 마감 보고',4.4,3.67,1.85,0.2,11,'ink',align=PP_ALIGN.CENTER,margin=0)
line(s,3.55,3.54,4.05,3.54,'red',2,dash=True)
txt(s,'확인 필요',2.75,4.48,1.5,0.28,13,'red',True,align=PP_ALIGN.CENTER,margin=0)
txt(s,'두 수치를 모두 인용하되, 기준 통일 전까지 의사결정 자료에 출처를 병기합니다.',1.15,5.12,5.45,0.45,12,'steel',margin=0)
# risk list
for i,(h,d,co) in enumerate([('수익성 편차','채널별 기여이익 기준 필요','amber'),('재고 회전 편차','물량 배분 기준 필요','navy'),('데이터 정합성','영업·재무 매출 기준 통일','red')]):
    y=2.22+i*1.18
    rect(s,7.55,y,4.65,0.88,'white',line='mist',radius=True)
    rect(s,7.78,y+0.23,0.14,0.34,co)
    txt(s,h,8.18,y+0.16,1.55,0.22,14,'navy',True,margin=0)
    txt(s,d,9.75,y+0.19,2.12,0.32,12,'ink',margin=0)
notes(s,'리스크는 세 가지입니다. 첫째 채널별 수익성 편차, 둘째 재고 회전 편차, 셋째 데이터 정합성입니다. 특히 매출은 영업팀 주간 보고 기준 누적 3억 원, 재무팀 마감 기준 누적 3.4억 원으로 불일치가 있습니다. 발표에서는 두 수치를 모두 인용하고, 확인 필요 과제로 남깁니다. 다음 분기 집행 전에 매출 기준과 채널별 기여이익, 재고 회전 지표를 통일해야 합니다.')

# 8 close / decision
s=prs.slides.add_slide(blank); base(s,dark=True); title(s,'마무리 및 의사결정 요청','강한 초기 수요를, 더 높은 효율의 성장으로 전환합니다.',dark=True)
rect(s,0.78,2.2,7.1,2.9,'white',radius=True)
txt(s,'오늘 요청드리는 결정',1.15,2.55,3.1,0.28,14,'amber',True,margin=0)
add_bullet(s,'2차 캠페인을 고효율 채널 중심으로 재배분',1.15,3.05,5.8,0.38,17,'navy','amber')
add_bullet(s,'다음 분기 추가 마케팅 예산 8억 원 승인',1.15,3.72,5.8,0.38,17,'navy','amber')
add_bullet(s,'매출 기준·수익성·재고 회전 확인 과제 병행',1.15,4.39,5.8,0.38,17,'navy','amber')
rect(s,8.55,2.2,3.85,2.9,'amber',radius=True)
txt(s,'다음 분기 목표',8.95,2.62,2.95,0.25,14,'navy',True,align=PP_ALIGN.CENTER,margin=0)
txt(s,'효율 중심\n확장',8.95,3.18,2.95,0.82,32,'navy',True,align=PP_ALIGN.CENTER,margin=0)
txt(s,'수요는 확인됐다.\n이제 운영 기준을 세운다.',8.95,4.35,2.95,0.45,13,'navy',align=PP_ALIGN.CENTER,margin=0)
txt(s,'감사합니다',0.8,6.15,5.0,0.35,18,'white',True,margin=0)
notes(s,'마무리하겠습니다. Aurora Kit은 6주 만에 12,400대 판매, 목표 대비 18% 초과라는 강한 초기 수요를 확인했습니다. 이제 필요한 결정은 2차 캠페인을 고효율 채널 중심으로 재배분하는 것, 다음 분기 추가 마케팅 예산 8억 원을 승인하는 것, 그리고 매출 기준·수익성·재고 회전 확인 과제를 함께 운영하는 것입니다. 수요는 확인됐고, 다음 분기는 효율 중심 확장으로 전환하겠습니다.')

prs.save(OUT)
print(f'created {OUT.resolve()} with {len(prs.slides)} slides')
