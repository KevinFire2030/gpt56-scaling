from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pathlib import Path

OUT = Path('deck.pptx')
W, H = Inches(13.333333), Inches(7.5)

NAVY = '1B2A4A'; AMBER = 'F5A623'; PAPER = 'FAF7F0'; INK = '22222A'; STEEL = '7C93B8'
WHITE = 'FFFFFF'; LIGHT_STEEL = 'E7ECF4'; PALE_AMBER = 'FFF0D4'; PALE_NAVY = 'DDE4F0'; MUTED = '5E6472'; GREEN = '2E7D6E'; RED = 'B04A4A'
FONT = '맑은 고딕'

def rgb(h): return RGBColor.from_string(h)

def set_fill(shape, color, transparency=0):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color); shape.fill.transparency = transparency

def set_line(shape, color=None, width=1, transparency=0):
    if color is None:
        shape.line.fill.background(); return
    shape.line.color.rgb = rgb(color); shape.line.width = Pt(width); shape.line.transparency = transparency

def add_box(slide, x, y, w, h, fill=WHITE, line=None, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    set_fill(s, fill); set_line(s, line)
    if radius:
        try: s.adjustments[0] = 0.08
        except Exception: pass
    return s

def set_text(shape, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT,
             valign=MSO_ANCHOR.MIDDLE, margin=0.08, font=FONT, line_spacing=1.05):
    tf = shape.text_frame; tf.clear()
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = text
    # PowerPoint can split newline-delimited text into multiple paragraphs/runs.
    # Apply formatting to every generated paragraph so subsequent lines do not
    # fall back to the theme color (especially important on dark slides).
    for para in tf.paragraphs:
        para.alignment = align; para.line_spacing = line_spacing
        for r in para.runs:
            r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return shape

def add_text(slide, text, x, y, w, h, size=20, color=INK, bold=False,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE, margin=0.03, name=None):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name: s.name = name
    return set_text(s, text, size, color, bold, align, valign, margin)

def add_rich(slide, runs, x, y, w, h, size=20, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.MIDDLE):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = s.text_frame; tf.clear(); tf.margin_left=tf.margin_right=Inches(.03); tf.margin_top=tf.margin_bottom=0
    tf.vertical_anchor=valign
    p=tf.paragraphs[0]; p.alignment=align
    for text,color,bold,sz in runs:
        r=p.add_run(); r.text=text; r.font.name=FONT; r.font.size=Pt(sz or size); r.font.bold=bold; r.font.color.rgb=rgb(color)
    return s

def add_title(slide, eyebrow, title, subtitle=None, dark=False):
    tc = WHITE if dark else NAVY; sc = PALE_NAVY if dark else MUTED
    add_text(slide, eyebrow.upper(), .65, .28, 5.5, .3, 10, AMBER, True, name='Eyebrow')
    # Use the real title placeholder for machine-parseable titles.
    ph = slide.shapes.title
    ph.left, ph.top, ph.width, ph.height = Inches(.65), Inches(.62), Inches(12.0), Inches(.58)
    set_text(ph, title, 25, tc, True, valign=MSO_ANCHOR.MIDDLE, margin=0)
    ph.name = 'Title'
    if subtitle: add_text(slide, subtitle, .67, 1.20, 11.8, .36, 12, sc, False)
    line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.65), Inches(1.55), Inches(12.0), Inches(.025)); set_fill(line, STEEL if not dark else WHITE, 55); set_line(line,None)

def add_footer(slide, num, source='출처: brief.md | 출시 후 6주'):
    add_text(slide, source, .67, 7.13, 5.6, .18, 8.5, MUTED)
    add_text(slide, f'{num:02d}', 12.18, 7.08, .48, .23, 9, STEEL, True, PP_ALIGN.RIGHT)

def add_note(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text

def metric_card(slide, x, y, w, h, value, label, accent=AMBER, detail=None):
    add_box(slide,x,y,w,h,WHITE,LIGHT_STEEL)
    add_text(slide,value,x+.22,y+.18,w-.44,.45,25,accent,True)
    add_text(slide,label,x+.22,y+.68,w-.44,.30,12,NAVY,True)
    if detail: add_text(slide,detail,x+.22,y+1.02,w-.44,h-1.12,9.5,MUTED)

def pill(slide,text,x,y,w,fill,color=WHITE):
    s=add_box(slide,x,y,w,.34,fill,None,True); set_text(s,text,9.5,color,True,PP_ALIGN.CENTER,margin=.02)

prs=Presentation(); prs.slide_width=W; prs.slide_height=H
# use title-only layout for all slides so every slide has a genuine title placeholder
layout=prs.slide_layouts[5]

# 1 Cover
s=prs.slides.add_slide(layout); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=rgb(NAVY)
# decorative geometry
shape=s.shapes.add_shape(MSO_SHAPE.ARC, Inches(9.15), Inches(-.55), Inches(5.3), Inches(5.3)); shape.fill.background(); shape.line.color.rgb=rgb(AMBER); shape.line.width=Pt(13); shape.line.transparency=12
shape.rotation=18
for i,(x,y,wh,c,t) in enumerate([(10.55,3.75,1.35,AMBER,8),(11.65,4.95,.55,STEEL,12),(9.55,5.35,.85,WHITE,78)]):
    q=s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(wh), Inches(wh)); set_fill(q,c,t); set_line(q,None)
pill(s,'EXECUTIVE REVIEW',.68,.52,1.95,AMBER,NAVY)
ph=s.shapes.title; ph.left=Inches(.68); ph.top=Inches(1.58); ph.width=Inches(8.8); ph.height=Inches(1.12); set_text(ph,'Aurora Kit\n출시 6주 성과 보고',31,WHITE,True,valign=MSO_ANCHOR.TOP,margin=0); ph.name='Title'
add_text(s,'강한 초기 수요를 수익성 있는 성장으로 전환',.72,3.10,7.7,.45,17,PALE_NAVY,False)
add_text(s,'경영진 회의  |  10분',.72,5.82,4.5,.32,12,WHITE,True)
add_text(s,'의사결정 요청: 2차 캠페인 예산 8억 원 및 채널 재배분 원칙',.72,6.25,8.8,.42,12,AMBER,True)
add_footer(s,1)
add_note(s,'오늘은 Aurora Kit 출시 후 6주 성과와 다음 분기 실행 결정을 보고드립니다. 핵심은 초기 수요는 기대 이상이지만 채널별 수익성과 재고 회전 편차를 줄여야 한다는 점입니다. 10분 안에 성과, 리스크, 요청 의사결정을 순서대로 말씀드리겠습니다. (약 45초)')

# 2 Agenda
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER)
add_title(s,'01 / ROADMAP','오늘의 논의는 세 가지입니다','성과 확인 → 리스크 정리 → 실행 승인')
items=[('01','출시 성과','목표 대비 판매와 핵심 퍼널'),('02','편차와 리스크','채널·재고·매출 수치 확인'),('03','다음 분기 결정','8억 원과 재배분 원칙')]
for i,(n,t,d) in enumerate(items):
    x=.72+i*4.12
    add_box(s,x,2.05,3.72,3.85,WHITE,LIGHT_STEEL)
    add_text(s,n,x+.25,2.30,.68,.50,24,AMBER,True)
    add_text(s,t,x+.25,3.15,3.0,.43,19,NAVY,True)
    add_text(s,d,x+.25,3.80,3.03,.85,12,MUTED,False,valign=MSO_ANCHOR.TOP)
    if i<2:
        add_text(s,'→',x+3.68,3.42,.42,.42,19,STEEL,True,PP_ALIGN.CENTER)
add_box(s,.72,6.20,11.95,.53,NAVY,None)
set_text(s.shapes[-1],'결론: 고효율 채널 중심으로 2차 캠페인을 재배분',14,WHITE,True,PP_ALIGN.CENTER)
add_footer(s,2)
add_note(s,'보고는 세 부분입니다. 먼저 판매와 퍼널 성과를 확인하고, 채널 편차와 매출 수치 불일치를 리스크로 정리합니다. 마지막으로 다음 분기 8억 원 예산과 고효율 채널 중심 재배분 원칙을 요청드립니다. (약 45초)')

# 3 Executive headline
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
add_title(s,'02 / EXECUTIVE SUMMARY','초기 수요는 강합니다. 이제 효율을 선택할 때입니다')
metric_card(s,.72,2.05,3.55,2.05,'12,400대','6주 누적 판매량',AMBER,'핵심 수요 지표')
metric_card(s,4.48,2.05,3.55,2.05,'+18%','목표 대비 판매 초과',NAVY,'계획을 상회한 초기 반응')
metric_card(s,8.24,2.05,3.55,2.05,'4.2%','웹 방문→구매 전환율',STEEL,'디지털 퍼널 기준')
add_box(s,.72,4.48,11.07,1.45,PALE_AMBER,None)
add_text(s,'KEY TAKEAWAY',.98,4.75,1.85,.25,9.5,NAVY,True)
add_text(s,'성장 신호는 확인됐지만, 채널별 수익성과 재고 회전 편차가 다음 성장을 제한합니다.',.98,5.08,10.25,.48,17,NAVY,True)
add_footer(s,3)
add_note(s,'6주 누적 판매는 12,400대로 목표를 18% 초과했습니다. 웹사이트 방문 대비 구매 전환율은 4.2%입니다. 초기 시장 반응은 분명히 강합니다. 다만 브리프가 지적하듯 채널별 수익성과 재고 회전 편차가 커, 다음 캠페인은 총량 확대보다 효율 선택이 우선입니다. (약 1분)')

# 4 Scorecard
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER)
add_title(s,'03 / PERFORMANCE','핵심 지표는 성장과 품질을 함께 지지합니다','브리프에 제공된 6주 기준 수치')
rows=[('판매량','12,400대','목표 대비 +18%','상회',AMBER),('구매 전환율','4.2%','웹사이트 방문 대비','관찰',STEEL),('반품률','1.9%','출시 6주 기준','낮은 마찰',GREEN),('고객 조사','61개 기업','재구매 의향 조사 응답','신호 확보',NAVY)]
add_text(s,'지표',.86,1.82,2.25,.3,10,MUTED,True); add_text(s,'실적',3.12,1.82,2.15,.3,10,MUTED,True); add_text(s,'기준',5.55,1.82,3.65,.3,10,MUTED,True); add_text(s,'판단',10.22,1.82,1.52,.3,10,MUTED,True)
for i,(a,b,c,d,col) in enumerate(rows):
    y=2.25+i*1.03
    add_box(s,.72,y,11.45,.82,WHITE,LIGHT_STEEL,False)
    add_text(s,a,.94,y+.14,2.00,.4,13,NAVY,True)
    add_text(s,b,3.12,y+.12,2.10,.44,18,col,True)
    add_text(s,c,5.55,y+.15,3.92,.38,11,MUTED)
    pill(s,d,10.15,y+.23,1.48,col,WHITE)
add_text(s,'해석',.74,6.55,.54,.26,9,AMBER,True)
add_text(s,'판매 초과와 낮은 반품률은 긍정적입니다. 고객 반응의 방향성은 61개 기업 응답의 상세 결과 확인이 필요합니다.',1.34,6.47,10.83,.42,11,INK)
add_footer(s,4)
add_note(s,'지표를 함께 보면 판매량은 목표 대비 18% 상회했고, 구매 전환율은 4.2%, 반품률은 1.9%입니다. 재구매 의향 조사에는 61개 기업 고객이 응답했습니다. 단, 브리프에는 재구매 의향의 긍정 비율이 없어 방향성을 단정하지 않았습니다. 상세 응답 분석을 후속 과제로 두겠습니다. (약 1분 10초)')

# 5 Channel + customer
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
add_title(s,'04 / CHANNEL & CUSTOMER','채널은 “확대”보다 “선별”이 우선입니다','채널별 실제 수치 미제공 — 아래는 의사결정 프레임 제안')
# matrix
add_text(s,'수익성 ↑',.78,1.90,1.0,.28,10,NAVY,True)
add_text(s,'재고 회전 →',5.17,5.90,1.05,.28,10,NAVY,True)
add_box(s,1.32,2.12,4.50,1.72,PALE_AMBER,None,False); add_box(s,5.87,2.12,4.50,1.72,NAVY,None,False)
add_box(s,1.32,3.89,4.50,1.72,LIGHT_STEEL,None,False); add_box(s,5.87,3.89,4.50,1.72,PALE_NAVY,None,False)
add_text(s,'유지·개선',1.65,2.45,3.82,.35,16,NAVY,True)
add_text(s,'수익성 높음 / 회전 낮음',1.65,2.92,3.82,.32,11,MUTED)
add_text(s,'우선 확대',6.20,2.45,3.82,.35,17,WHITE,True)
add_text(s,'수익성 높음 / 회전 높음',6.20,2.92,3.82,.32,11,PALE_NAVY)
add_text(s,'축소 검토',1.65,4.22,3.82,.35,16,NAVY,True)
add_text(s,'수익성 낮음 / 회전 낮음',1.65,4.69,3.82,.32,11,MUTED)
add_text(s,'조건부 운영',6.20,4.22,3.82,.35,16,NAVY,True)
add_text(s,'수익성 낮음 / 회전 높음',6.20,4.69,3.82,.32,11,MUTED)
# side customer card
add_box(s,10.62,2.12,2.05,3.49,PAPER,LIGHT_STEEL)
add_text(s,'61',10.87,2.55,1.52,.65,30,AMBER,True,PP_ALIGN.CENTER)
add_text(s,'개 기업 고객',10.87,3.20,1.52,.35,12,NAVY,True,PP_ALIGN.CENTER)
add_text(s,'재구매 의향\n조사 응답',10.92,3.78,1.42,.70,11,MUTED,False,PP_ALIGN.CENTER)
add_text(s,'상세 긍정률\n확인 필요',10.92,4.75,1.42,.48,9.5,RED,True,PP_ALIGN.CENTER)
add_text(s,'가정',.74,6.53,.52,.25,9,AMBER,True)
add_text(s,'실제 채널별 수치가 제공되지 않아 채널명·순위를 임의로 만들지 않았습니다. 다음 집행 전 4분면 분류를 완료합니다.',1.32,6.43,11.0,.45,10.5,INK)
add_footer(s,5)
add_note(s,'브리프에는 채널별 수익성과 재고 회전 편차가 크다고 명시되어 있지만 채널별 실제 수치는 없습니다. 따라서 채널명이나 순위를 임의로 만들지 않고, 수익성과 재고 회전의 2축 프레임으로 의사결정 기준을 제안합니다. 두 지표가 모두 높은 채널을 우선 확대하고 나머지는 개선·조건부 운영·축소 검토로 분류합니다. 고객 조사 61개 기업의 상세 긍정률도 확인이 필요합니다. (약 1분 15초)')

# 6 Revenue discrepancy
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER)
add_title(s,'05 / RISK','누적 매출 0.4억 원 차이는 “검증 필요”입니다','영업 주간 보고와 재무 마감 기준이 불일치')
add_box(s,.72,2.05,4.58,2.60,WHITE,LIGHT_STEEL)
pill(s,'영업팀 주간 보고',.98,2.36,1.72,STEEL,WHITE)
add_text(s,'3억 원',.98,3.05,3.85,.72,31,NAVY,True)
add_text(s,'운영 시점 기준 누적 매출',.98,3.85,3.85,.35,11,MUTED)
add_box(s,7.33,2.05,4.58,2.60,WHITE,LIGHT_STEEL)
pill(s,'재무팀 마감',7.59,2.36,1.48,NAVY,WHITE)
add_text(s,'3.4억 원',7.59,3.05,3.85,.72,31,NAVY,True)
add_text(s,'마감 기준 누적 매출',7.59,3.85,3.85,.35,11,MUTED)
add_text(s,'≠',5.67,2.92,1.00,.65,30,AMBER,True,PP_ALIGN.CENTER)
add_box(s,4.98,4.86,2.64,.83,AMBER,None)
set_text(s.shapes[-1],'차이 0.4억 원',17,NAVY,True,PP_ALIGN.CENTER)
add_box(s,.72,6.05,11.19,.62,NAVY,None)
set_text(s.shapes[-1],'확인 과제  |  집계 시점 · 취소/반품 반영 · 매출 인식 기준 · 데이터 오너',12.5,WHITE,True,PP_ALIGN.CENTER)
add_footer(s,6)
add_note(s,'매출은 영업팀 주간 보고 3억 원과 재무팀 마감 3.4억 원으로 불일치합니다. 두 수치를 모두 보고드리며, 차이 0.4억 원은 검증이 필요합니다. 집계 시점, 취소·반품 반영, 매출 인식 기준을 대조하고 재무를 최종 데이터 오너로 지정하는 안을 제안합니다. 이 검증 전에는 채널 ROI를 확정 수치로 보고하지 않겠습니다. (약 1분 15초)')

# 7 Action plan
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(WHITE)
add_title(s,'06 / NEXT QUARTER','8억 원은 성과 확인 후 단계적으로 집행합니다','제안 가정: 70% 확대 / 20% 실험 / 10% 보류')
alloc=[('70%','5.6억 원','고효율 채널 확대',NAVY,6.78),('20%','1.6억 원','신규·개선 실험',STEEL,1.94),('10%','0.8억 원','성과 연동 보류',AMBER,.97)]
x=.74
for pct,amt,lab,col,width in alloc:
    add_box(s,x,2.13,width,1.05,col,None,False)
    add_text(s,pct,x+.12,2.25,width-.24,.30,16,WHITE if col!=AMBER else NAVY,True,PP_ALIGN.CENTER)
    add_text(s,amt,x+.12,2.61,width-.24,.28,10,WHITE if col!=AMBER else NAVY,True,PP_ALIGN.CENTER)
    x+=width+.06
for i,(n,t,d) in enumerate([('1','분류','채널별 수익성·재고 회전 4분면 확정'),('2','집행','고효율군에 70% 우선 배정'),('3','게이트','주간 전환·반품·회전 리뷰 후 확대')]):
    x=.74+i*4.05
    add_box(s,x,3.70,3.70,2.12,PAPER,LIGHT_STEEL)
    pill(s,n,x+.24,3.95,.42,AMBER,NAVY)
    add_text(s,t,x+.79,3.90,2.30,.35,15,NAVY,True)
    add_text(s,d,x+.24,4.56,3.10,.78,11,MUTED,False,valign=MSO_ANCHOR.TOP)
add_text(s,'가정',.74,6.32,.52,.25,9,AMBER,True)
add_text(s,'배분 비율은 브리프 미제공 사항으로, 리스크를 통제하기 위한 제안입니다. 승인 후 실제 채널 데이터로 조정합니다.',1.32,6.22,11.0,.46,10.5,INK)
add_footer(s,7)
add_note(s,'다음 분기 추가 마케팅 예산 요청액은 8억 원입니다. 집행 비율은 브리프에 없으므로 리스크 통제를 위한 제안 가정으로 70%, 20%, 10%를 제시합니다. 5.6억 원은 고효율 채널 확대, 1.6억 원은 신규·개선 실험, 0.8억 원은 성과 연동 보류입니다. 채널 분류를 먼저 확정하고 주간 게이트를 통과할 때만 확대합니다. (약 1분 20초)')

# 8 Decision ask
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(PAPER)
add_title(s,'07 / DECISION','오늘 세 가지 결정을 요청드립니다')
asks=[('01','예산','다음 분기 추가 마케팅\n예산 8억 원 승인','APPROVE'),('02','원칙','고효율 채널 중심\n재배분 원칙 승인','ALIGN'),('03','통제','매출 검증 오너와\n주간 집행 게이트 지정','ASSIGN')]
for i,(n,t,d,tag) in enumerate(asks):
    x=.72+i*4.10
    add_box(s,x,2.05,3.70,3.95,WHITE,LIGHT_STEEL)
    add_text(s,n,x+.25,2.31,.56,.42,18,AMBER,True)
    pill(s,tag,x+2.52,2.34,.85,NAVY,WHITE)
    add_text(s,t,x+.25,3.05,3.03,.38,17,NAVY,True)
    add_text(s,d,x+.25,3.76,3.05,.95,13,INK,True,valign=MSO_ANCHOR.TOP)
    add_text(s,'결정 시 즉시 실행',x+.25,5.22,3.02,.30,10,GREEN,True)
add_box(s,.72,6.31,11.90,.48,AMBER,None)
set_text(s.shapes[-1],'강한 수요를 확인했습니다. 다음 단계는 “더 많이”가 아니라 “더 효율적으로”입니다.',13,NAVY,True,PP_ALIGN.CENTER)
add_footer(s,8)
add_note(s,'오늘 요청드리는 결정은 세 가지입니다. 첫째, 다음 분기 추가 마케팅 예산 8억 원 승인. 둘째, 고효율 채널 중심 재배분 원칙 승인. 셋째, 누적 매출 불일치 검증 오너와 주간 집행 게이트 지정입니다. 승인되면 채널별 4분면 분류와 매출 대사를 먼저 완료하고 단계 집행을 시작하겠습니다. (약 1분)')

# 9 Close
s=prs.slides.add_slide(layout); s.background.fill.solid(); s.background.fill.fore_color.rgb=rgb(NAVY)
add_title(s,'08 / CLOSE','수요는 입증됐습니다.\n이제 수익성 있는 성장으로 전환하겠습니다.',dark=True)
add_box(s,.72,2.30,11.90,2.26,WHITE,None)
add_rich(s,[('12,400대',AMBER,True,28),('  ·  ',STEEL,False,24),('+18%',NAVY,True,28),('  ·  ',STEEL,False,24),('8억 원',NAVY,True,28)],1.10,2.67,11.12,.68,align=PP_ALIGN.CENTER)
add_text(s,'6주 누적 판매  ·  목표 초과  ·  다음 분기 요청 예산',1.10,3.50,11.12,.36,12,MUTED,False,PP_ALIGN.CENTER)
add_text(s,'의사결정 후 첫 실행: 매출 수치 검증 → 채널 4분면 확정 → 단계 집행',1.04,5.25,11.26,.42,15,WHITE,True,PP_ALIGN.CENTER)
pill(s,'THANK YOU',5.62,6.18,2.10,AMBER,NAVY)
add_footer(s,9)
add_note(s,'정리하겠습니다. Aurora Kit는 6주 누적 12,400대로 목표를 18% 초과하며 수요를 입증했습니다. 다음 분기에는 8억 원을 일괄 집행하지 않고, 매출 수치 검증과 채널 분류를 거쳐 고효율 채널 중심으로 단계 집행하겠습니다. 세 가지 의사결정을 부탁드립니다. 감사합니다. (약 40초)')

# core metadata
prs.core_properties.title='Aurora Kit 출시 6주 성과 보고'
prs.core_properties.subject='경영진 10분 보고 — 성과, 리스크, 다음 분기 의사결정'
prs.core_properties.author='경영진 보고 자료'
prs.core_properties.comments='brief.md 기반. 채널별 실제 수치와 예산 배분 비율의 가정은 슬라이드에 명시.'
prs.save(OUT)
print(f'created {OUT.resolve()} with {len(prs.slides)} slides')
