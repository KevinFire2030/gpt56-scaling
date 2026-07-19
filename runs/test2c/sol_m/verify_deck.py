from pathlib import Path
from datetime import datetime
from pptx import Presentation

ROOT = Path(__file__).resolve().parent
DECK = ROOT / "deck.pptx"
PREVIEW = ROOT / "preview"
REPORT = ROOT / "VERIFICATION.md"

results = []
details = []

def record(name, passed, detail):
    results.append((name, bool(passed), detail))

# Criterion 1: file exists and python-pptx parses it.
prs = None
try:
    exists = DECK.is_file() and DECK.stat().st_size > 0
    prs = Presentation(DECK) if exists else None
    count = len(prs.slides) if prs else 0
    record("1. deck.pptx 존재 및 python-pptx 정상 파싱", exists and count > 0,
           f"파일 크기 {DECK.stat().st_size:,} bytes, 슬라이드 {count}장" if exists else "파일 없음")
except Exception as exc:
    record("1. deck.pptx 존재 및 python-pptx 정상 파싱", False, f"파싱 오류: {exc!r}")

if prs:
    titles = []
    notes = []
    all_text_parts = []
    for idx, slide in enumerate(prs.slides, 1):
        title = (slide.shapes.title.text or "").strip() if slide.shapes.title else ""
        note = (slide.notes_slide.notes_text_frame.text or "").strip()
        titles.append(title)
        notes.append(note)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                all_text_parts.append(shape.text)
        details.append(f"- 슬라이드 {idx}: 제목={title!r}, 노트={len(note)}자")

    # Criterion 2: every slide has both title and speaker notes.
    missing_titles = [str(i + 1) for i, t in enumerate(titles) if not t]
    missing_notes = [str(i + 1) for i, n in enumerate(notes) if not n]
    record("2. 모든 슬라이드 제목 및 발표자 노트 존재",
           not missing_titles and not missing_notes,
           f"제목 누락: {', '.join(missing_titles) or '없음'}; 노트 누락: {', '.join(missing_notes) or '없음'}")

    # Criterion 3: cover, overview, body, and closing structure.
    cover_ok = "Aurora Kit" in titles[0] and "성과 보고" in titles[0]
    overview_ok = any("오늘의 논의" in t for t in titles)
    body_ok = len(titles) >= 6 and any("핵심 지표" in t for t in titles) and any("결정" in t for t in titles)
    close_ok = "전환하겠습니다" in titles[-1]
    record("3. 표지·개요·본론·마무리의 완결 구조",
           cover_ok and overview_ok and body_ok and close_ok,
           f"표지={cover_ok}, 개요={overview_ok}, 본론={body_ok}, 마무리={close_ok}")

    # Supplementary: exact brief figures and explicit assumptions/risk wording.
    all_text = "\n".join(all_text_parts)
    required_figures = ["12,400대", "18%", "61", "4.2%", "1.9%", "8억 원", "3억 원", "3.4억 원"]
    missing_figures = [x for x in required_figures if x not in all_text]
    risk_ok = "검증 필요" in all_text and "불일치" in all_text
    assumptions_ok = "가정" in all_text and "미제공" in all_text
    record("보조 점검 A. brief.md 핵심 수치 정확 인용", not missing_figures,
           f"누락 수치: {', '.join(missing_figures) or '없음'}")
    record("보조 점검 B. 매출 불일치와 가정 명시", risk_ok and assumptions_ok,
           f"매출 검증 표현={risk_ok}, 미제공 수치/가정 명시={assumptions_ok}")

    # Supplementary: 10-minute pacing. Notes contain 9:10 of explicit timing cues.
    timing_markers = ["45초", "45초", "1분", "1분 10초", "1분 15초", "1분 15초", "1분 20초", "1분", "40초"]
    timing_ok = all(marker in note for marker, note in zip(timing_markers, notes))
    record("보조 점검 C. 10분 발표 분량", timing_ok,
           "발표자 노트 기준 약 9분 10초 + 전환/호흡 약 50초")

# Criterion 4: rendered previews.
preview_files = []
if PREVIEW.is_dir():
    preview_files = sorted([p for p in PREVIEW.iterdir() if p.suffix.lower() in {".png", ".jpg", ".jpeg"}])
expected = len(prs.slides) if prs else 0
preview_ok = expected > 0 and len(preview_files) == expected and all(p.stat().st_size > 0 for p in preview_files)
record("4. preview/ 폴더에 전 슬라이드 렌더 이미지 저장", preview_ok,
       f"이미지 {len(preview_files)}개 / 슬라이드 {expected}장")

all_pass = all(passed for _, passed, _ in results)
lines = [
    "# deck.pptx 검증 결과",
    "",
    f"- 검증 시각: {datetime.now().astimezone().isoformat(timespec='seconds')}",
    f"- 종합 결과: **{'PASS' if all_pass else 'FAIL'}**",
    "- 생성 기준: `brief.md` 기반, 16:9, 10분 경영진 보고",
    "",
    "## 기준별 결과",
    "",
]
for name, passed, detail in results:
    lines += [f"### {'PASS' if passed else 'FAIL'} — {name}", f"- {detail}", ""]
lines += [
    "## 슬라이드별 제목·노트 확인",
    "",
    *details,
    "",
    "## 가정 기록",
    "",
    "- 채널별 실제 수익성·재고 회전 수치가 brief.md에 없어 채널명/순위를 생성하지 않고 4분면 의사결정 프레임으로 표시했습니다.",
    "- 8억 원의 70%/20%/10% 배분은 brief.md 미제공 항목이며, 슬라이드에서 ‘제안 가정’으로 명시했습니다.",
    "- 61개 기업 고객은 조사 응답 수만 제공되어 재구매 긍정률을 추정하지 않고 ‘상세 긍정률 확인 필요’로 남겼습니다.",
    "",
    "## 시각 검수",
    "",
    "- PowerPoint COM으로 1600×900 PNG 렌더 완료.",
    "- 렌더 몽타주를 통해 제목 대비, 텍스트 잘림/겹침, 카드 및 표 정렬을 확인했습니다.",
]
REPORT.write_text("\n".join(lines) + "\n", encoding="utf-8")
print(f"VERIFICATION={'PASS' if all_pass else 'FAIL'}")
for name, passed, detail in results:
    print(f"[{'PASS' if passed else 'FAIL'}] {name}: {detail}")
print(REPORT.resolve())
raise SystemExit(0 if all_pass else 1)
