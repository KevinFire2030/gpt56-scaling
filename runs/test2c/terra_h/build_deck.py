from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
from pathlib import Path

OUT = Path('deck.pptx')

# Aurora Kit brand tokens from brief.md
NAVY = '1B2A4A'; AMBER = 'F5A623'; PAPER = 'FAF7F0'; INK = '22222A'; STEEL = '7C93B8'; WHITE = 'FFFFFF'; LIGHT = 'E9EDF4'; MUTED = '5C6678'

def rgb(hexv):
    return RGBColor.from_string(hexv)

def fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()

def line(shape, color, width=1):
    shape.line.color.rgb = rgb(color); shape.line.width = Pt(width)

def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, font='Aptos', align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.04):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = Inches(margin); tf.margin_right = Inches(margin); tf.margin_top = Inches(margin); tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return tb

def add_para(tf, text, size=18, color=INK, bold=False, level=0, space_after=8):
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.level = level; p.space_after = Pt(space_after)
    r = p.add_run(); r.text = text; r.font.name = 'Aptos'; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return p

def rect(slide,x,y,w,h,color, radius=False, linecolor=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    fill(s,color)
    if linecolor: line(s,linecolor)
    return s

def title(slide, kicker, heading, sub=None, page=None):
    add_text(slide, kicker.upper(), 0.62, 0.35, 7.7, 0.27, 10, AMBER, True)
    add_text(slide, heading, 0.62, 0.68, 11.25, 0.58, 27, NAVY, True)
    if sub: add_text(slide, sub, 0.64, 1.28, 11.1, 0.34, 12, MUTED)
    rect(slide,0.62,1.72,1.05,0.05,AMBER)
    if page: add_text(slide, f'{page:02d}', 12.18, 7.0, 0.45, 0.25, 10, STEEL, True, align=PP_ALIGN.RIGHT)

def footer(slide):
    rect(slide,0.62,7.15,12.05,0.015,LIGHT)
    add_text(slide, 'AURORA KIT  |  출시 6주 성과 보고', 0.62,7.20,4.5,0.18,8,STEEL,True)

def notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text

def metric(slide, x,y,w,h, value, label, accent=AMBER, note=None):
    rect(slide,x,y,w,h,WHITE,True,LIGHT)
    rect(slide,x,y,0.07,h,accent)
    add_text(slide,value,x+0.25,y+0.25,w-0.45,0.55,28,NAVY,True)
    add_text(slide,label,x+0.25,y+0.95,w-0.4,0.35,12,INK,True)
    if note: add_text(slide,note,x+0.25,y+1.35,w-0.4,0.3,10,MUTED)

def bullets(slide, items, x,y,w,h, size=17):
    box=slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=Inches(0.05); tf.margin_right=Inches(0.05)
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.level=0; p.space_after=Pt(13)
        p.font.name='Aptos'; p.font.size=Pt(size); p.font.color.rgb=rgb(INK); p.font.bold=False
        p._p.get_or_add_pPr().insert(0, OxmlElement('a:buChar')) if False else None
        # bullet shown with a text glyph to maintain reliable rendering
        p.text='•  '+item
    return box

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
blank=prs.slide_layouts[6]

def slide():
    s=prs.slides.add_slide(blank); bg=s.background.fill; bg.solid(); bg.fore_color.rgb=rgb(PAPER); return s

# 1 Cover
s=slide(); rect(s,0,0,13.333,7.5,NAVY); rect(s,0.65,0.72,1.15,0.07,AMBER)
add_text(s,'AURORA KIT',0.7,1.10,5.0,0.4,14,AMBER,True)
add_text(s,'출시 6주 성과 보고',0.66,1.65,8.8,0.78,34,WHITE,True)
add_text(s,'초기 수요는 강력합니다. 다음 분기는 효율 중심으로 재배분합니다.',0.70,2.62,8.7,0.4,17,'D8E1F0')
rect(s,0.68,4.35,4.8,1.35,'24385F',True)
add_text(s,'경영진 회의  |  다음 분기 실행 결정 요청',0.95,4.64,4.3,0.27,13,WHITE,True)
add_text(s,'보고 기준: 출시 후 6주',0.95,5.05,4.1,0.22,11,'D8E1F0')
add_text(s,'CONFIDENTIAL  ·  MANAGEMENT REPORT',0.7,6.88,4.4,0.18,9,'A7B8D3',True)
notes(s,'[약 40초] Aurora Kit 출시 후 6주 성과를 보고드립니다. 핵심은 초기 수요가 예상보다 강하다는 점입니다. 다만 채널별 수익성과 재고 회전의 편차가 커서, 오늘은 다음 분기 마케팅 예산 8억 원의 재배분 방향에 대한 결정을 요청드리겠습니다.')

#2 overview
s=slide(); title(s,'01  |  개요','오늘의 결론과 의사결정', '10분: 성과 → 고객·채널 시사점 → 리스크 → 요청사항',2)
steps=[('01','출시 성과','판매·전환·품질의 핵심 지표'),('02','무엇을 배웠나','고객 반응과 채널 운영 시사점'),('03','무엇을 결정할까','고효율 채널 중심의 예산 재배분')]
for i,(n,h,d) in enumerate(steps):
    x=0.8+i*4.12; rect(s,x,2.25,3.58,2.35,WHITE,True,LIGHT); add_text(s,n,x+0.28,2.55,0.55,0.27,13,AMBER,True); add_text(s,h,x+0.28,3.05,2.9,0.38,20,NAVY,True); add_text(s,d,x+0.28,3.65,2.95,0.48,13,MUTED)
add_text(s,'결론  |  수요 모멘텀은 유지하되, 투자 기준을 “확장”에서 “효율·회전”으로 전환합니다.',0.82,5.35,11.6,0.42,18,NAVY,True)
footer(s); notes(s,'[약 45초] 발표는 세 부분으로 진행합니다. 먼저 핵심 성과를 숫자로 확인하고, 이어 고객과 채널에서 얻은 시사점을 보겠습니다. 마지막으로 실행 리스크를 짚고, 다음 분기에 대한 의사결정 요청을 드리겠습니다.')

#3 scorecard
s=slide(); title(s,'02  |  출시 성과','핵심 지표는 초기 수요의 강도를 보여줍니다','목표 초과 판매, 전환, 반품률을 한 화면에서 확인',3)
metric(s,0.72,2.12,2.9,1.85,'12,400대','출시 후 6주 누적 판매량')
metric(s,3.84,2.12,2.9,1.85,'18%','목표 대비 판매 초과율')
metric(s,6.96,2.12,2.9,1.85,'4.2%','웹사이트 방문 대비 구매 전환율')
metric(s,10.08,2.12,2.55,1.85,'1.9%','반품률',STEEL)
rect(s,0.72,4.55,11.91,1.05,'EAF0F8',True)
add_text(s,'경영진 시사점',1.02,4.81,1.75,0.25,12,AMBER,True)
add_text(s,'판매가 목표를 18% 상회한 가운데, 전환과 반품률은 출시 직후 수요가 실제 구매와 제품 경험으로 연결되고 있음을 보여줍니다.',1.02,5.16,10.95,0.30,15,NAVY,True)
footer(s); notes(s,'[약 75초] 6주 누적 판매량은 12,400대로 목표를 18% 초과했습니다. 웹사이트 방문 대비 구매 전환율은 4.2%, 반품률은 1.9%입니다. 이 네 지표를 함께 보면 초기 관심이 실제 구매로 이어지고 있으며, 제품 경험의 초기 품질 신호도 안정적이라는 판단입니다.')

#4 Channel
s=slide(); title(s,'03  |  채널 시사점','다음 캠페인의 기준은 “규모”보다 “효율”입니다','세부 채널 수익성·재고 회전은 운영 데이터로 즉시 점검 필요',4)
rect(s,0.75,2.05,3.65,3.6,WHITE,True,LIGHT); add_text(s,'관찰',1.05,2.35,2.6,0.27,13,AMBER,True); add_text(s,'초기 수요는 예상보다 강함',1.05,2.85,2.9,0.7,23,NAVY,True); add_text(s,'6주 판매 12,400대\n목표 대비 +18%',1.05,4.15,2.5,0.6,16,INK,True)
rect(s,4.83,2.05,3.65,3.6,WHITE,True,LIGHT); add_text(s,'과제',5.13,2.35,2.6,0.27,13,AMBER,True); add_text(s,'채널별 수익성·재고 회전 편차',5.13,2.85,2.95,0.7,23,NAVY,True); add_text(s,'브리프의 핵심 리스크\n채널별 운영 지표 확인 필요',5.13,4.15,2.85,0.6,15,INK,True)
rect(s,8.91,2.05,3.65,3.6,NAVY,True); add_text(s,'실행 원칙',9.21,2.35,2.6,0.27,13,AMBER,True); add_text(s,'고효율 채널로 재배분',9.21,2.85,2.9,0.7,23,WHITE,True); add_text(s,'2차 캠페인은 효율과 회전을\n주간 단위로 관리',9.21,4.15,2.85,0.6,15,'D8E1F0',True)
add_text(s,'가정: 채널별 상세 수치는 본 브리프에 포함되지 않아, 본 슬라이드는 방향성과 관리 기준만 제시합니다.',0.78,6.20,11.7,0.28,10,MUTED)
footer(s); notes(s,'[약 75초] 채널별 상세 수익성과 재고 회전 수치는 현재 브리프에 포함되어 있지 않습니다. 따라서 이를 새로 추정하지 않고, 의사결정의 원칙을 분명히 하겠습니다. 2차 캠페인은 판매 규모 자체보다 채널별 수익성과 회전이 검증된 곳에 투자하고, 주간 단위로 기준을 재점검해야 합니다.')

#5 customer
s=slide(); title(s,'04  |  고객 반응','기업 고객의 재구매 의향을 다음 검증의 출발점으로','정성 반응은 유지하되, 표본 기반의 후속 확인이 필요',5)
metric(s,0.78,2.05,3.35,2.0,'61개','재구매 의향 조사 응답\n기업 고객',AMBER)
rect(s,4.55,2.05,7.9,2.0,WHITE,True,LIGHT); add_text(s,'의미',4.9,2.36,1.0,0.25,13,AMBER,True); add_text(s,'기업 고객 반응을 다음 분기 반복 구매·확장 제안의 우선 검증 대상으로 활용',4.9,2.82,6.9,0.55,21,NAVY,True)
rect(s,0.78,4.62,11.67,1.05,'EAF0F8',True); add_text(s,'다음 확인',1.08,4.90,1.45,0.25,12,AMBER,True); add_text(s,'응답 61개 기업 고객을 기반으로 재구매 의향의 구체적 전환 조건과 채널별 반영 방법을 확인합니다.',2.62,4.88,9.25,0.30,15,NAVY,True)
footer(s); notes(s,'[약 60초] 기업 고객 대상 재구매 의향 조사에는 61개 고객이 응답했습니다. 이 숫자만으로 재구매 규모를 단정하지는 않겠습니다. 대신 다음 분기에는 이 응답 기반으로 실제 재구매 전환 조건을 확인하고, 검증된 조건을 고효율 채널 운영에 연결하겠습니다.')

#6 operating pulse
s=slide(); title(s,'05  |  운영 상태','성장과 품질을 동시에 관리할 수 있는 초기 신호','수요 확대 전, 회전·수익성 편차를 통제하는 운영 체계가 핵심',6)
# simple visual
rect(s,0.78,2.1,5.4,3.55,WHITE,True,LIGHT); add_text(s,'제품 경험 신호',1.10,2.42,2.4,0.28,13,AMBER,True); add_text(s,'반품률 1.9%',1.10,2.95,3.8,0.50,29,NAVY,True); add_text(s,'초기 수요 확대 과정에서\n제품 경험을 계속 추적',1.10,3.75,3.6,0.55,16,INK)
rect(s,6.60,2.1,5.95,3.55,NAVY,True); add_text(s,'운영 우선순위',6.95,2.42,2.7,0.28,13,AMBER,True); bullets(s,['채널별 수익성 기준 설정','재고 회전 주간 모니터링','2차 캠페인 투자 재배분'],6.93,2.96,4.85,1.65,16)
# make bullet white manually
for shp in s.shapes:
    if hasattr(shp,'text_frame') and shp.left>=Inches(6.9) and shp.top>=Inches(2.9) and shp.top<Inches(5):
        for p in shp.text_frame.paragraphs:
            for r in p.runs: r.font.color.rgb=rgb(WHITE)
footer(s); notes(s,'[약 60초] 반품률은 1.9%로, 현재 시점에서는 제품 경험 측면의 초기 모니터링 신호로 보겠습니다. 다만 수요 확대를 바로 넓히기보다, 채널별 수익성 기준과 재고 회전을 주간 단위로 관리해야 합니다. 성과를 유지하는 운영의 초점은 이 세 가지입니다.')

#7 revenue reconciliation
s=slide(); title(s,'06  |  리스크 및 확인 필요','누적 매출 기준이 두 보고에서 불일치합니다','의사결정 전 단일 기준·정산 시점을 확정해야 합니다',7)
metric(s,0.85,2.12,4.7,2.05,'3억 원','영업팀 주간 보고 기준\n누적 매출',AMBER)
metric(s,7.78,2.12,4.7,2.05,'3.4억 원','재무팀 마감 기준\n누적 매출',STEEL)
# arrow/issue
add_text(s,'≠',6.12,2.67,1.0,0.6,35,AMBER,True,align=PP_ALIGN.CENTER)
rect(s,0.85,4.75,11.63,0.95,'FFF2D8',True); add_text(s,'확인 필요',1.15,5.03,1.25,0.23,12,AMBER,True); add_text(s,'영업팀 주간 보고와 재무팀 마감 보고의 누적 매출 수치가 불일치 — 다음 분기 예산 집행 전 단일 기준으로 검증합니다.',2.57,5.00,9.15,0.33,15,NAVY,True)
footer(s); notes(s,'[약 75초] 중요한 확인 사항입니다. 영업팀 주간 보고 기준 누적 매출은 3억 원이고, 재무팀 마감 기준은 3.4억 원입니다. 두 수치는 불일치합니다. 오늘 이 차이를 해석하거나 하나를 선택하지 않겠습니다. 다음 분기 예산 집행 전에 정산 기준과 시점을 단일화해 검증하는 과제로 명확히 남기겠습니다.')

#8 decision
s=slide(); title(s,'07  |  의사결정 요청','다음 분기: 8억 원을 고효율 채널에 재배분','수요 모멘텀을 유지하면서 수익성과 회전의 편차를 줄입니다',8)
rect(s,0.8,2.05,3.3,3.65,NAVY,True); add_text(s,'요청',1.15,2.42,1.5,0.28,13,AMBER,True); add_text(s,'추가 마케팅\n예산 승인',1.15,2.95,2.4,0.75,24,WHITE,True); add_text(s,'8억 원',1.15,4.35,2.2,0.45,28,AMBER,True)
items=[('1','재배분','고효율 채널 우선'),('2','운영','재고 회전 주간 관리'),('3','통제','매출 기준 단일화·검증')]
for i,(n,h,d) in enumerate(items):
    y=2.08+i*1.16; rect(s,4.62,y,7.78,0.91,WHITE,True,LIGHT); add_text(s,n,4.93,y+0.25,0.35,0.25,14,AMBER,True); add_text(s,h,5.55,y+0.20,1.3,0.26,16,NAVY,True); add_text(s,d,7.35,y+0.21,4.2,0.25,15,INK)
footer(s); notes(s,'[약 75초] 의사결정 요청은 다음 분기 추가 마케팅 예산 8억 원입니다. 단순한 확장 예산이 아니라, 고효율 채널에 재배분하고 재고 회전을 주간 관리하는 조건부 실행안입니다. 동시에 매출 기준의 불일치를 해소해 투자 판단의 기반을 정리하겠습니다.')

#9 close
s=slide(); rect(s,0,0,13.333,7.5,NAVY); rect(s,0.68,0.75,1.25,0.07,AMBER)
add_text(s,'CLOSING',0.70,1.12,2.0,0.28,12,AMBER,True)
add_text(s,'강한 초기 수요를\n효율적인 성장으로 전환하겠습니다.',0.68,1.70,9.2,1.52,31,WHITE,True)
rect(s,0.70,4.45,11.95,1.08,'24385F',True); add_text(s,'오늘의 요청  |  8억 원 예산 승인 · 고효율 채널 재배분 · 매출 기준 단일화 검증',1.04,4.82,11.1,0.30,16,'D8E1F0',True)
add_text(s,'Aurora Kit  |  출시 6주 성과 보고',0.70,6.85,4.2,0.2,9,'A7B8D3',True)
notes(s,'[약 35초] 정리하겠습니다. Aurora Kit은 출시 6주 만에 목표를 상회하는 수요를 확인했습니다. 다음 단계는 더 넓게 쓰는 것이 아니라 더 효율적으로 쓰는 것입니다. 8억 원의 추가 마케팅 예산을 고효율 채널에 재배분하고, 회전과 매출 기준을 통제하는 실행안을 승인해 주시길 요청드립니다.')

# Basic document properties
prs.core_properties.title='Aurora Kit | 출시 6주 성과 보고'
prs.core_properties.subject='경영진 회의용 신제품 출시 성과 및 다음 분기 의사결정 요청'
prs.core_properties.author='Aurora Kit Team'
prs.save(OUT)
print(f'Created {OUT.resolve()} with {len(prs.slides)} slides')
