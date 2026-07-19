from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pathlib import Path

OUT = Path('deck.pptx')
W, H = Inches(13.333), Inches(7.5)
NAVY = '1B2A4A'; AMBER = 'F5A623'; PAPER = 'FAF7F0'; INK = '22222A'; STEEL = '7C93B8'; WHITE = 'FFFFFF'; PALE = 'E9EEF5'; LIGHT_AMBER='FFF1D7'; MUTED='64748B'

prs = Presentation(); prs.slide_width=W; prs.slide_height=H
blank = prs.slide_layouts[6]

def rgb(hex):
    hex=hex.lstrip('#'); return RGBColor.from_string(hex)

def rect(slide,x,y,w,h,fill, line=None, radius=False):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb=rgb(fill)
    shp.line.color.rgb=rgb(line or fill)
    if radius:
        shp.adjustments[0]=0.12
    return shp

def text(slide,content,x,y,w,h,size=18,color=INK,bold=False,font='Aptos',align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, margin=0.05):
    box=slide.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=box.text_frame; tf.clear(); tf.word_wrap=True
    tf.margin_left=tf.margin_right=Inches(margin); tf.margin_top=tf.margin_bottom=Inches(margin); tf.vertical_anchor=valign
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=content
    r.font.name=font; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=rgb(color)
    return box

def title(slide,num,heading,kicker=None):
    text(slide, f'{num:02d}', .55,.35,.5,.3,10,AMBER,True)
    text(slide, heading,.55,.68,12.1,.55,27,NAVY,True)
    if kicker: text(slide,kicker,.57,1.28,11.8,.28,11,MUTED)
    rect(slide,.55,1.65,12.15,.025,PALE)

def footer(slide,n):
    text(slide,'AURORA KIT  |  출시 6주 성과',.57,7.08,4,.18,8,STEEL,True)
    text(slide,f'{n} / 09',11.75,7.08,.9,.18,8,STEEL,True,align=PP_ALIGN.RIGHT)

def add_notes(slide, note):
    tf=slide.notes_slide.notes_text_frame
    tf.text = note

def bullet(slide, content,x,y,w, accent=AMBER, size=17):
    text(slide,'•',x,y,.18,.28,size,accent,True)
    text(slide,content,x+.24,y,w-.24,.42,size,INK)

def metric(slide,x,y,w,h,value,label,sub=None,fill=WHITE,accent=NAVY):
    rect(slide,x,y,w,h,fill,PALE,True)
    rect(slide,x,y,.08,h,accent,accent,True)
    text(slide,value,x+.28,y+.23,w-.48,.42,25,accent,True)
    text(slide,label,x+.28,y+.78,w-.48,.32,12,INK,True)
    if sub: text(slide,sub,x+.28,y+1.13,w-.48,.25,10,MUTED)

# 1 cover
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); rect(s,0,0,.23,7.5,AMBER)
text(s,'AURORA KIT',.75,.72,3.1,.35,14,AMBER,True)
text(s,'신제품 출시 성과 보고',.75,1.32,8.6,.7,34,NAVY,True)
text(s,'출시 6주 성과와 다음 분기 실행 결정',.77,2.13,7.5,.4,18,INK)
rect(s,.77,3.0,2.15,.05,AMBER)
text(s,'경영진 회의  |  10분',.77,6.45,3,.25,12,STEEL,True)
# visual bars
for i,(h,c) in enumerate([(1.0,STEEL),(1.7,NAVY),(2.35,AMBER),(2.65,NAVY),(3.05,AMBER),(3.55,NAVY)]): rect(s,9.4+i*.42,5.95-h,.25,h,c,c,True)
text(s,'6 WEEKS',9.35,6.3,2.9,.35,13,NAVY,True,align=PP_ALIGN.CENTER)
add_notes(s,'오늘은 Aurora Kit 출시 후 6주 성과를 간결히 보고드리고, 다음 분기 2차 캠페인 예산 8억 원의 재배분 결정을 요청드리겠습니다. 핵심은 초기 수요는 강했지만 채널별 수익성과 재고 회전 편차에 대응해야 한다는 점입니다.')

# 2 agenda
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,2,'오늘의 논의: 성과 확인에서 실행 결정까지','10분 | 성과 4분 · 시사점/리스크 3분 · 의사결정 3분')
items=[('01','출시 성과 요약','초기 수요와 핵심 KPI'),('02','채널·고객 시사점','효율의 차이와 고객 신호'),('03','리스크와 확인 과제','매출 기준 불일치, 재고 회전'),('04','다음 분기 결정 요청','고효율 채널 중심 2차 캠페인')]
for i,(n,h,d) in enumerate(items):
    x=.7+(i%2)*6.15; y=2.05+(i//2)*2.05
    rect(s,x,y,5.65,1.55,WHITE,PALE,True); text(s,n,x+.28,y+.25,.55,.28,13,AMBER,True); text(s,h,x+.28,y+.62,4.8,.35,19,NAVY,True); text(s,d,x+.28,y+1.08,4.95,.25,11,MUTED)
footer(s,2); add_notes(s,'구성은 네 부분입니다. 먼저 숫자로 출시 성과를 확인하고, 채널과 고객 반응에서 의미를 읽겠습니다. 이어 매출 기준과 재고의 확인 과제를 짚은 뒤, 다음 분기 실행 결정을 요청드리겠습니다.')

# 3 executive summary
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,3,'출시 6주: 수요는 강하고, 실행은 선택과 집중이 필요합니다','핵심 메시지')
metric(s,.72,2.05,2.8,1.65,'12,400대','누적 판매량','출시 후 6주',WHITE,AMBER)
metric(s,3.75,2.05,2.8,1.65,'+18%','목표 대비 판매 초과율','예상 상회',WHITE,NAVY)
metric(s,6.78,2.05,2.8,1.65,'4.2%','웹 방문 대비 구매 전환율','디지털 구매 신호',WHITE,NAVY)
metric(s,9.81,2.05,2.8,1.65,'1.9%','반품률','제품 경험 안정성',WHITE,STEEL)
rect(s,.72,4.22,11.89,1.5,NAVY,NAVY,True)
text(s,'결론',1.0,4.54,1.0,.25,12,AMBER,True)
text(s,'2차 캠페인은 고효율 채널 중심으로 재배분하고, 재고 회전 및 매출 기준을 즉시 정렬합니다.',1.0,4.93,10.7,.38,21,WHITE,True)
footer(s,3); add_notes(s,'출시 6주의 결론입니다. 12,400대를 판매하며 목표를 18% 초과했습니다. 웹 전환율은 4.2%, 반품률은 1.9%입니다. 제품 수요와 경험은 긍정적입니다. 다만 다음 단계의 성과는 예산을 모든 채널에 넓게 쓰기보다 고효율 채널에 집중하는 실행력에 달려 있습니다.')

#4 KPI
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,4,'목표를 18% 초과 달성하며 초기 수요를 확인했습니다','핵심 지표와 목표 대비 성과')
# bar chart sales index
text(s,'판매 목표 대비 지수',.8,2.0,2.5,.25,13,NAVY,True)
rect(s,1.0,3.05,8.55,.52,PALE,PALE,True); rect(s,1.0,3.05,7.25,.52,NAVY,NAVY,True); rect(s,8.25,2.82,.04,.98,AMBER,AMBER)
text(s,'목표 100',8.0,3.83,.7,.2,10,MUTED,align=PP_ALIGN.CENTER); text(s,'실적 118',8.72,3.18,1.0,.25,14,AMBER,True)
text(s,'출시 6주 누적 판매량',.8,4.45,3,.28,13,NAVY,True); text(s,'12,400대',.8,4.85,3,.55,34,NAVY,True)
rect(s,6.9,4.35,5.55,1.24,LIGHT_AMBER,LIGHT_AMBER,True); text(s,'해석',7.2,4.62,.7,.25,12,AMBER,True); text(s,'초기 수요는 계획보다 강함',7.2,4.96,4.8,.27,17,INK,True)
text(s,'※ 브리프에 제공된 판매 초과율을 기준으로 표시',.8,6.4,5.5,.2,10,MUTED)
footer(s,4); add_notes(s,'목표 대비 성과를 보면 판매는 18% 초과했습니다. 누적 판매량은 12,400대입니다. 이 수치는 시장 진입과 초기 수요 검증이 성공적으로 시작됐다는 신호입니다. 이 강한 출발을 다음 분기에 더 효율적인 성장으로 전환하는 것이 핵심입니다.')

#5 channel
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,5,'채널별로 수익성과 재고 회전 편차가 커졌습니다','채널별 시사점 | 제공 데이터에 채널별 수치가 없어, 방향성 중심으로 제시')
# flow cards
for x,h,d,c in [(0.75,'고효율 채널','2차 캠페인 예산을 우선 배분',NAVY),(4.65,'개선 필요 채널','수익성·재고 회전 원인 점검',AMBER),(8.55,'공통 운영','주간 관리 지표와 판단 기준 통일',STEEL)]:
 rect(s,x,2.35,3.35,2.32,WHITE,PALE,True); rect(s,x,2.35,3.35,.12,c,c,True); text(s,h,x+.3,2.78,2.8,.35,19,NAVY,True); text(s,d,x+.3,3.45,2.75,.55,14,INK)
text(s,'권고: “전체 확대”가 아닌 “효율 증명 후 확장” 원칙으로 2차 캠페인을 설계',.85,5.5,11.5,.4,20,NAVY,True,align=PP_ALIGN.CENTER)
footer(s,5); add_notes(s,'브리프는 채널별 수익성과 재고 회전의 편차가 크다고 명시합니다. 세부 채널별 수치는 제공되지 않았으므로 이 슬라이드는 방향성 중심의 권고입니다. 고효율 채널에는 예산을 우선 배분하고, 개선 필요 채널은 수익성과 재고 회전의 원인을 먼저 점검해야 합니다. 공통으로는 주간 지표와 판단 기준을 맞추겠습니다.')

#6 customer
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,6,'고객 반응은 구매 전환과 낮은 반품률로 확인됩니다','고객 반응')
metric(s,.85,2.05,3.5,1.72,'61개','재구매 의향 조사 응답','기업 고객 응답',WHITE,NAVY)
metric(s,4.92,2.05,3.5,1.72,'4.2%','웹 방문 대비 구매 전환율','구매 전환 신호',WHITE,AMBER)
metric(s,8.99,2.05,3.5,1.72,'1.9%','반품률','제품 경험 안정성',WHITE,STEEL)
rect(s,.85,4.42,11.64,1.22,WHITE,PALE,True)
text(s,'고객 관점의 함의',1.15,4.68,2.2,.25,12,AMBER,True)
text(s,'전환을 확장하되, 재구매 의향 응답을 다음 캠페인의 타깃·메시지 설계에 활용합니다.',1.15,5.02,10.3,.3,17,INK,True)
footer(s,6); add_notes(s,'고객 관점에서 세 가지 신호가 있습니다. 기업 고객 61개가 재구매 의향 조사에 응답했고, 웹 전환율은 4.2%, 반품률은 1.9%입니다. 이 결과는 제품 수용성과 구매 여정 모두에 긍정적 신호를 줍니다. 다음 캠페인에서는 이 61개 응답에서 얻는 재구매 요인을 타깃과 메시지에 반영하겠습니다.')

#7 risks
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,7,'성과를 확장하기 전에 두 가지 운영 리스크를 해소해야 합니다','리스크와 확인 필요 사항')
rect(s,.78,2.05,5.68,2.7,WHITE,PALE,True); text(s,'01  매출 기준 불일치',1.08,2.42,3.8,.3,19,NAVY,True); text(s,'영업팀 주간 보고',1.08,3.06,2.2,.25,12,MUTED); text(s,'3억 원',1.08,3.32,1.6,.44,27,NAVY,True); text(s,'재무팀 마감 보고',3.72,3.06,2.2,.25,12,MUTED); text(s,'3.4억 원',3.72,3.32,1.7,.44,27,AMBER,True)
text(s,'확인 필요: 집계 기준·시점 정렬 후 단일 기준 확정',1.08,4.2,4.9,.25,12,INK,True)
rect(s,6.85,2.05,5.68,2.7,WHITE,PALE,True); text(s,'02  재고 회전 편차',7.15,2.42,3.8,.3,19,NAVY,True); text(s,'채널별 수익성과 재고 회전의 편차가 확대',7.15,3.08,4.7,.45,16,INK); text(s,'확인 필요: 고효율 채널 우선 배분 전, 채널별 회전·마진 점검',7.15,4.02,4.8,.38,12,INK,True)
text(s,'원칙: 수치의 신뢰도를 먼저 정렬하고, 그 위에서 예산을 재배분합니다.',.9,5.65,11.5,.32,18,NAVY,True,align=PP_ALIGN.CENTER)
footer(s,7); add_notes(s,'확장 전에 명확히 관리해야 할 리스크입니다. 첫째, 누적 매출이 영업팀 주간 보고에서는 3억 원, 재무팀 마감 기준에서는 3.4억 원으로 불일치합니다. 이는 확인 필요 사항이며, 집계 기준과 시점을 정렬해 단일 기준을 확정하겠습니다. 둘째, 채널별 수익성과 재고 회전 편차입니다. 예산 재배분 전 반드시 채널별 회전과 마진을 점검하겠습니다.')

#8 decision
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,PAPER); title(s,8,'다음 분기: 8억 원을 고효율 채널 중심으로 재배분해 주십시오','의사결정 요청')
rect(s,.8,2.0,3.25,3.38,NAVY,NAVY,True); text(s,'요청 예산',1.1,2.42,1.6,.25,13,AMBER,True); text(s,'8억 원',1.1,2.8,2.2,.55,33,WHITE,True); text(s,'다음 분기 추가\n마케팅 예산',1.1,3.56,2,.55,16,WHITE)
steps=[('1','고효율 채널 우선','수익성·회전이 검증된 채널에 집중'),('2','저효율 채널 개선','원인 점검 후 조건부 확장'),('3','주간 운영 체계','매출 기준·재고 회전 단일 대시보드')]
for i,(n,h,d) in enumerate(steps):
 y=2.0+i*1.18; rect(s,4.65,y,7.75,.92,WHITE,PALE,True); rect(s,4.9,y+.2,.42,.42,AMBER,AMBER,True); text(s,n,4.9,y+.25,.42,.18,11,NAVY,True,align=PP_ALIGN.CENTER); text(s,h,5.55,y+.18,2.7,.25,16,NAVY,True); text(s,d,8.1,y+.2,3.8,.25,12,INK)
footer(s,8); add_notes(s,'요청드리는 결정은 다음 분기 추가 마케팅 예산 8억 원의 승인과 운용 원칙입니다. 첫째 고효율 채널에 우선 배분합니다. 둘째 저효율 채널은 원인 점검과 개선 조건을 충족할 때 확장합니다. 셋째 매출 기준과 재고 회전을 하나의 주간 운영 체계로 관리하겠습니다. 이를 통해 강한 초기 수요를 효율적인 성장으로 연결하겠습니다.')

#9 close
s=prs.slides.add_slide(blank); rect(s,0,0,13.333,7.5,NAVY); rect(s,0,0,.23,7.5,AMBER)
text(s,'AURORA KIT',.75,.8,2.5,.3,13,AMBER,True)
text(s,'강한 초기 수요를\n선택과 집중의 성장으로.',.75,1.45,9.0,1.4,34,WHITE,True)
rect(s,.78,3.32,1.7,.05,AMBER)
text(s,'오늘의 결정',.78,4.0,1.2,.25,12,AMBER,True)
text(s,'다음 분기 8억 원 승인 · 고효율 채널 우선 재배분 · 데이터 기준 정렬',.78,4.4,11.4,.34,19,WHITE,True)
text(s,'감사합니다',.78,6.5,2,.28,14,STEEL,True)
add_notes(s,'마무리하겠습니다. Aurora Kit은 출시 6주 동안 강한 초기 수요를 확인했습니다. 이제는 고효율 채널에 집중하고, 매출 기준과 재고 회전의 운영 신뢰도를 정렬해야 합니다. 다음 분기 8억 원 예산 승인과 이 실행 원칙에 대한 결정을 요청드립니다. 감사합니다.')

prs.save(OUT)
print(OUT.resolve())
print('slides', len(prs.slides))
