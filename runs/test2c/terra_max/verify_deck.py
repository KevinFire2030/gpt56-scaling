from pathlib import Path
from pptx import Presentation

DECK=Path('deck.pptx'); PREVIEW=Path('preview')
expected_titles = [
    '신제품 출시 성과 보고',
    '오늘의 논의: 성과 확인에서 실행 결정까지',
    '출시 6주: 수요는 강하고, 실행은 선택과 집중이 필요합니다',
    '목표를 18% 초과 달성하며 초기 수요를 확인했습니다',
    '채널별로 수익성과 재고 회전 편차가 커졌습니다',
    '고객 반응은 구매 전환과 낮은 반품률로 확인됩니다',
    '성과를 확장하기 전에 두 가지 운영 리스크를 해소해야 합니다',
    '다음 분기: 8억 원을 고효율 채널 중심으로 재배분해 주십시오',
    '강한 초기 수요를\n선택과 집중의 성장으로.'
]
checks=[]
try:
    prs=Presentation(str(DECK))
    checks.append(('1. deck.pptx 파일 존재 및 python-pptx 정상 파싱', DECK.exists() and len(prs.slides)==9, f'파일 존재={DECK.exists()}, 파싱 슬라이드 수={len(prs.slides)}'))
except Exception as e:
    prs=None; checks.append(('1. deck.pptx 파일 존재 및 python-pptx 정상 파싱', False, repr(e)))

if prs:
    title_ok=True; note_ok=True; detail=[]
    for i,(slide,expected) in enumerate(zip(prs.slides, expected_titles),1):
        all_text='\n'.join(sh.text for sh in slide.shapes if getattr(sh,'has_text_frame',False))
        note=slide.notes_slide.notes_text_frame.text.strip()
        found=expected in all_text
        has_note=bool(note)
        title_ok &= found; note_ok &= has_note
        detail.append(f'- 슬라이드 {i}: 제목={"OK" if found else "FAIL"}, 노트={"OK" if has_note else "FAIL"} ({len(note)}자)')
    checks.append(('2. 모든 슬라이드에 제목과 발표자 노트 존재', title_ok and note_ok, '\n'.join(detail)))
    # Cover is first, closing is last, agenda at second, body at 3-8
    structure_ok = len(prs.slides)==9 and expected_titles[0] in '\n'.join(sh.text for sh in prs.slides[0].shapes if getattr(sh,'has_text_frame',False)) and expected_titles[-1] in '\n'.join(sh.text for sh in prs.slides[-1].shapes if getattr(sh,'has_text_frame',False))
    checks.append(('3. 표지·개요·본론·마무리의 완결 구조', structure_ok, '1=표지, 2=개요, 3–8=본론, 9=마무리; 총 9장(10분 발표용)'))
else:
    checks.extend([('2. 모든 슬라이드에 제목과 발표자 노트 존재',False,'PPTX 파싱 실패'),('3. 표지·개요·본론·마무리의 완결 구조',False,'PPTX 파싱 실패')])

contact=PREVIEW/'contact_sheet.png'
# On Windows glob is case-insensitive; exclude the separately generated contact sheet.
pngs=sorted(p for p in PREVIEW.glob('*.png') if p.name.lower() != 'contact_sheet.png')
preview_ok=len(pngs)==9 and all(p.stat().st_size>0 for p in pngs) and contact.exists() and contact.stat().st_size>0
checks.append(('4. 슬라이드 이미지 렌더 및 preview/ 저장', preview_ok, f'PowerPoint PNG={len(pngs)}장, contact_sheet.png={contact.exists()}'))

lines=['# VERIFICATION', '', '검증 대상: `deck.pptx`', '검증 일시: 생성 직후 자동 검증', '']
for name,ok,detail in checks:
    lines += [f'## [{"PASS" if ok else "FAIL"}] {name}', detail, '']
lines += ['## 제작 가정', '- 채널별 세부 수치가 브리프에 제공되지 않아, 해당 슬라이드는 수치를 추정하지 않고 브리프의 방향성(고효율 채널 우선·수익성/재고 회전 점검)만 제시했습니다.', '- 누적 매출은 영업팀 주간 보고 3억 원과 재무팀 마감 보고 3.4억 원을 모두 정확히 표기하고, “불일치” 및 확인 과제로 명시했습니다.', '']
Path('VERIFICATION.md').write_text('\n'.join(lines), encoding='utf-8')
print('\n'.join(lines))
if not all(ok for _,ok,_ in checks): raise SystemExit(1)
